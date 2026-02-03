#!/usr/bin/env python3
"""
Script para gerar apresentação PowerPoint do SIGESC
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

def set_shape_gradient(shape, color1, color2):
    """Define gradiente para um shape"""
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = 45
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def add_title_slide(prs, title, subtitle=""):
    """Adiciona slide de título"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background gradient shape
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RgbColor(15, 23, 42)  # slate-900
    bg.line.fill.background()
    
    # Decorative circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7), Inches(-1), Inches(5), Inches(5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RgbColor(30, 64, 175)  # blue-800
    circle.fill.fore_color.brightness = 0.3
    circle.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RgbColor(148, 163, 184)  # slate-400
        p.alignment = PP_ALIGN.LEFT
    
    return slide

def add_content_slide(prs, title, content_items, accent_color=RgbColor(59, 130, 246)):
    """Adiciona slide de conteúdo com bullets"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RgbColor(15, 23, 42)
    bg.line.fill.background()
    
    # Accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), prs.slide_height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        
        p.text = f"• {item}"
        p.font.size = Pt(22)
        p.font.color.rgb = RgbColor(226, 232, 240)  # slate-200
        p.space_after = Pt(12)
        p.level = 0
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Adiciona slide com duas colunas"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RgbColor(15, 23, 42)
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RgbColor(34, 197, 94)  # green-500
    
    # Left content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(4.5), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = f"✓ {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = RgbColor(203, 213, 225)
        p.space_after = Pt(8)
    
    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.5), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RgbColor(59, 130, 246)  # blue-500
    
    # Right content
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.7), Inches(4.5), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = f"→ {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = RgbColor(203, 213, 225)
        p.space_after = Pt(8)
    
    return slide

def add_stats_slide(prs):
    """Adiciona slide de estatísticas"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RgbColor(15, 23, 42)
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SIGESC em Números"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    # Stats boxes
    stats = [
        ("5.000+", "Alunos Gerenciados", RgbColor(59, 130, 246)),
        ("15+", "Escolas Atendidas", RgbColor(34, 197, 94)),
        ("99,9%", "Disponibilidade", RgbColor(168, 85, 247)),
        ("24/7", "Suporte Técnico", RgbColor(249, 115, 22)),
    ]
    
    start_x = 0.3
    box_width = 2.2
    spacing = 0.2
    
    for i, (value, label, color) in enumerate(stats):
        x = start_x + i * (box_width + spacing)
        
        # Box background
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2), Inches(box_width), Inches(2.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RgbColor(30, 41, 59)  # slate-800
        box.line.color.rgb = RgbColor(51, 65, 85)  # slate-700
        
        # Value
        value_box = slide.shapes.add_textbox(Inches(x), Inches(2.3), Inches(box_width), Inches(1))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = slide.shapes.add_textbox(Inches(x), Inches(3.5), Inches(box_width), Inches(0.8))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.color.rgb = RgbColor(148, 163, 184)
        p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Números em constante crescimento com a expansão da rede de escolas atendidas"
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = RgbColor(100, 116, 139)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_contact_slide(prs):
    """Adiciona slide de contato"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background gradient
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RgbColor(30, 64, 175)  # blue-800
    bg.line.fill.background()
    
    # Decorative element
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6), Inches(3), Inches(6), Inches(6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RgbColor(37, 99, 235)  # blue-600
    circle.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Vamos Conversar?"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(6), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Entre em contato para uma demonstração gratuita"
    p.font.size = Pt(22)
    p.font.color.rgb = RgbColor(191, 219, 254)  # blue-200
    
    # Contact info
    contacts = [
        ("📱 WhatsApp:", "(94) 98422-3453"),
        ("📧 E-mail:", "contato@aprenderdigital.top"),
        ("🌐 Site:", "aprenderdigital.top"),
    ]
    
    y_pos = 3.5
    for icon_text, value in contacts:
        # Icon/Label
        label_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(2.5), Inches(0.5))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon_text
        p.font.size = Pt(20)
        p.font.color.rgb = RgbColor(191, 219, 254)
        
        # Value
        value_box = slide.shapes.add_textbox(Inches(2.8), Inches(y_pos), Inches(4), Inches(0.5))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RgbColor(255, 255, 255)
        
        y_pos += 0.6
    
    # Footer - Aprender Digital
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(0.5))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Projeto de propriedade da Aprender Digital"
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(147, 197, 253)  # blue-300
    
    return slide

def create_presentation():
    """Cria a apresentação completa"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Capa
    add_title_slide(prs, "SIGESC", "Sistema Integrado de Gestão Escolar")
    
    # Slide 2: O que é o SIGESC
    add_content_slide(prs, "O que é o SIGESC?", [
        "Sistema completo para gestão da rede municipal de educação",
        "Desenvolvido especialmente para a realidade das escolas brasileiras",
        "Acesso pela internet de qualquer lugar, 24 horas por dia",
        "Interface moderna, intuitiva e fácil de usar",
        "Funciona até mesmo sem conexão com a internet (modo offline)",
        "Segurança de dados com backup automático na nuvem",
    ], RgbColor(59, 130, 246))
    
    # Slide 3: Estatísticas
    add_stats_slide(prs)
    
    # Slide 4: Funcionalidades - Gestão de Alunos
    add_content_slide(prs, "Gestão Completa de Alunos", [
        "Cadastro completo com todos os dados pessoais e documentos",
        "Histórico escolar automático com todas as movimentações",
        "Controle de matrículas, transferências e remanejamentos",
        "Pré-matrícula online para novos alunos (autoatendimento)",
        "Busca e filtros avançados por nome, turma, status e escola",
        "Geração de declarações, fichas e documentos em PDF",
        "Registro de atestados médicos com bloqueio de frequência",
    ], RgbColor(34, 197, 94))
    
    # Slide 5: Funcionalidades - Notas e Frequência
    add_content_slide(prs, "Notas e Frequência", [
        "Lançamento de notas por bimestre/trimestre/semestre",
        "Suporte a notas numéricas e conceituais (Educação Infantil)",
        "Controle de frequência diário com cálculo automático",
        "Boletins individuais gerados automaticamente em PDF",
        "Atas de resultados finais para impressão",
        "Ficha individual do aluno completa",
        "Relatórios gerenciais para acompanhamento",
    ], RgbColor(168, 85, 247))
    
    # Slide 6: Funcionalidades - Gestão Escolar
    add_content_slide(prs, "Gestão de Escolas e Turmas", [
        "Cadastro de múltiplas escolas e unidades",
        "Organização de turmas por nível de ensino e série",
        "Alocação de professores por disciplina",
        "Calendário escolar integrado com feriados e eventos",
        "Controle de servidores (professores, funcionários)",
        "Sistema de avisos e comunicação interna",
        "Dashboard com visão geral da rede de ensino",
    ], RgbColor(249, 115, 22))
    
    # Slide 7: Funcionalidades - Secretaria
    add_content_slide(prs, "Ferramentas para Secretaria", [
        "Geração de documentos oficiais em PDF (declarações, históricos)",
        "Controle de matrículas e transferências entre escolas",
        "Ações em lote para atualização de status de alunos",
        "Relatórios consolidados por escola e rede",
        "Auditoria de todas as alterações no sistema",
        "Notificações em tempo real de novas pré-matrículas",
        "Exportação de dados para planilhas Excel",
    ], RgbColor(236, 72, 153))
    
    # Slide 8: Diferenciais
    add_content_slide(prs, "Diferenciais do SIGESC", [
        "Funciona offline - continue trabalhando sem internet",
        "Interface responsiva - acesse pelo celular ou tablet",
        "Suporte técnico dedicado e treinamento incluído",
        "Atualizações constantes sem custo adicional",
        "Dados seguros na nuvem com backup automático",
        "Personalização com a identidade visual do município",
        "Custo-benefício superior às soluções do mercado",
    ], RgbColor(20, 184, 166))
    
    # Slide 9: O que está implementado vs Futuro
    add_two_column_slide(
        prs,
        "Roadmap do SIGESC",
        "✅ Já Implementado",
        [
            "Gestão completa de alunos",
            "Matrículas e transferências",
            "Lançamento de notas",
            "Controle de frequência",
            "Geração de boletins/PDFs",
            "Pré-matrícula online",
            "Calendário escolar",
            "Sistema de avisos",
            "Modo offline",
            "Atestados médicos",
        ],
        "🚀 Em Desenvolvimento",
        [
            "App mobile nativo",
            "Portal do responsável",
            "Integração com INEP/Censo",
            "Relatórios avançados BI",
            "Biblioteca digital",
            "Gestão de transporte escolar",
            "Controle de merenda",
            "Integração com PNLD",
            "Módulo financeiro",
            "Comunicação via WhatsApp",
        ]
    )
    
    # Slide 10: Benefícios
    add_content_slide(prs, "Benefícios para a Rede de Ensino", [
        "Redução de até 70% no tempo de processos administrativos",
        "Eliminação de papelada e arquivos físicos",
        "Dados centralizados e acessíveis em tempo real",
        "Maior controle e transparência na gestão",
        "Facilidade na tomada de decisões com dados atualizados",
        "Padronização de processos em todas as escolas",
        "Economia de recursos com digitalização de documentos",
    ], RgbColor(234, 179, 8))
    
    # Slide 11: Contato
    add_contact_slide(prs)
    
    # Salva a apresentação
    output_path = "/app/SIGESC_Apresentacao.pptx"
    prs.save(output_path)
    print(f"Apresentação salva em: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
